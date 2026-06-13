"""
extractor.py — Text extraction module for AI File Search Assistant (V2).

V2 changes over V1:
    - Class-based architecture: BaseExtractor + per-type extractor classes
    - OCR support for images via RapidOCR (graceful fallback if not installed)
    - Scanned PDF detection: pages with little/no extractable text are OCR'd
    - PowerPoint support (.pptx via python-pptx)
    - Excel support (.xlsx via openpyxl)
    - Clean HTML extraction via BeautifulSoup (strips scripts, styles)
    - File metadata extraction (name, path, size, dates)
    - Text chunking with configurable size and overlap

Image captioning is intentionally NOT included — OCR only (reads text in
images/scanned pages, does not describe image content).

Supported file types:
    Plain text  : .txt
    Markup      : .md .json .xml
    HTML        : .html .htm          (cleaned via BeautifulSoup)
    PDF         : .pdf                (text + OCR fallback for scanned pages)
    Word        : .docx
    PowerPoint  : .pptx
    Excel       : .xlsx
    Spreadsheet : .csv
    Code        : .py .js .ts .java .c .cpp .h .hpp .cs .go .rs .php .rb .swift .kt
    Image       : .jpg .jpeg .png .bmp .gif .webp  (OCR via RapidOCR)

Public API:
    extract_text(file_path)                         → str   (V1 backward-compatible)
    extract_with_metadata(file_path, ...)           → ExtractionResult | None
    chunk_text(text, size, overlap)                 → list[str]
    get_supported_extensions()                      → set[str]

Install optional dependencies:
    pip install rapidocr-onnxruntime    # OCR for images and scanned PDFs
    pip install python-pptx             # PowerPoint
    pip install openpyxl                # Excel
    pip install beautifulsoup4          # Clean HTML extraction
"""

from __future__ import annotations

import csv
import logging
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

# ── Constants ─────────────────────────────────────────────────────────────────

MAX_TEXT_LENGTH: int = 50_000
DEFAULT_CHUNK_SIZE: int = 800
DEFAULT_CHUNK_OVERLAP: int = 150

# Pages with fewer characters than this are treated as scanned (→ OCR fallback)
SCANNED_PAGE_THRESHOLD: int = 100

# ── Logging ───────────────────────────────────────────────────────────────────

logger = logging.getLogger(__name__)


# ── Data classes ──────────────────────────────────────────────────────────────

@dataclass
class FileMetadata:
    """Filesystem metadata for an extracted file."""
    name: str
    path: str
    extension: str
    size: int
    created_date: str
    modified_date: str


@dataclass
class ExtractionResult:
    """Full V2 extraction output: cleaned content + metadata + text chunks."""
    content: str
    metadata: FileMetadata
    chunks: list[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        """Serialise to a plain dict (ready for SQLite or JSON storage)."""
        return {
            "path": self.metadata.path,
            "name": self.metadata.name,
            "extension": self.metadata.extension,
            "size": self.metadata.size,
            "created_date": self.metadata.created_date,
            "modified_date": self.metadata.modified_date,
            "content": self.content,
            "chunks": self.chunks,
        }


# ── Base extractor ────────────────────────────────────────────────────────────

class BaseExtractor(ABC):
    """Abstract base for all file-type extractors."""

    @abstractmethod
    def extract(self, path: Path) -> str:
        """Extract raw text from *path*. Returns empty string on failure."""
        ...

    @abstractmethod
    def supported_extensions(self) -> frozenset[str]:
        """Return the lowercase file extensions this extractor handles."""
        ...


# ── OCR utility ───────────────────────────────────────────────────────────────

def _run_ocr_on_bytes(image_bytes: bytes) -> str:
    """
    Run RapidOCR on raw image bytes.

    Returns extracted text, or an empty string when RapidOCR is not installed
    or the image yields no recognisable text.
    """
    try:
        from rapidocr_onnxruntime import RapidOCR  # noqa: PLC0415

        engine = RapidOCR()
        result, _ = engine(image_bytes)

        if not result:
            return ""

        return "\n".join(line[1] for line in result if line[1])

    except ImportError:
        logger.warning(
            "RapidOCR not installed — OCR unavailable. "
            "Install with: pip install rapidocr-onnxruntime"
        )
        return ""
    except Exception as exc:
        logger.error("OCR failed: %s", exc)
        return ""


# ── Extractor implementations ─────────────────────────────────────────────────

class TxtExtractor(BaseExtractor):
    """Reads plain text, markup, and source code files as UTF-8."""

    _EXTENSIONS = frozenset({
        ".txt", ".md", ".json", ".xml",
        ".py", ".js", ".ts", ".java", ".c", ".cpp",
        ".h", ".hpp", ".cs", ".go", ".rs", ".php",
        ".rb", ".swift", ".kt",
    })

    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def extract(self, path: Path) -> str:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception as exc:
            logger.error("TxtExtractor failed for '%s': %s", path.name, exc)
            return ""


class HTMLExtractor(BaseExtractor):
    """
    Extracts visible text from HTML files using BeautifulSoup.
    Falls back to raw file read if BeautifulSoup is not installed.
    """

    _EXTENSIONS = frozenset({".html", ".htm"})

    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def extract(self, path: Path) -> str:
        try:
            from bs4 import BeautifulSoup  # noqa: PLC0415

            raw = path.read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(raw, "html.parser")

            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()

            return soup.get_text(separator="\n")

        except ImportError:
            logger.warning(
                "BeautifulSoup not installed — returning raw HTML. "
                "Install with: pip install beautifulsoup4"
            )
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception as exc:
            logger.error("HTMLExtractor failed for '%s': %s", path.name, exc)
            return ""


class PDFExtractor(BaseExtractor):
    """
    Extracts text from PDFs using PyMuPDF.

    Scanned page detection: if a page yields fewer than
    SCANNED_PAGE_THRESHOLD characters, the page is rendered to an image
    and OCR is run on it instead. This makes scanned PDFs searchable
    without any manual configuration.
    """

    _EXTENSIONS = frozenset({".pdf"})

    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def extract(self, path: Path) -> str:
        try:
            import fitz  # PyMuPDF  # noqa: PLC0415

            doc = fitz.open(str(path))

            if doc.is_encrypted:
                logger.warning("Encrypted PDF, cannot extract: %s", path.name)
                doc.close()
                return ""

            pages: list[str] = []
            for page in doc:
                text = page.get_text()
                if len(text.strip()) < SCANNED_PAGE_THRESHOLD:
                    logger.info(
                        "Page %d of '%s' appears scanned — running OCR.",
                        page.number + 1, path.name,
                    )
                    text = self._ocr_page(page)
                pages.append(text)

            doc.close()
            return "\n".join(pages)

        except Exception as exc:
            logger.error("PDFExtractor failed for '%s': %s", path.name, exc)
            return ""

    @staticmethod
    def _ocr_page(page) -> str:
        """Render a fitz page to PNG bytes and pass it to the OCR utility."""
        try:
            pix = page.get_pixmap(dpi=150)
            return _run_ocr_on_bytes(pix.tobytes("png"))
        except Exception as exc:
            logger.warning("OCR fallback failed for page %d: %s", page.number + 1, exc)
            return ""


class DOCXExtractor(BaseExtractor):
    """Extracts all paragraph text from a Word document."""

    _EXTENSIONS = frozenset({".docx"})

    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def extract(self, path: Path) -> str:
        try:
            from docx import Document  # noqa: PLC0415

            doc = Document(str(path))
            return "\n".join(para.text for para in doc.paragraphs)

        except Exception as exc:
            logger.error("DOCXExtractor failed for '%s': %s", path.name, exc)
            return ""


class CSVExtractor(BaseExtractor):
    """Converts CSV rows to tab-separated text lines."""

    _EXTENSIONS = frozenset({".csv"})

    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def extract(self, path: Path) -> str:
        try:
            lines: list[str] = []
            with path.open(encoding="utf-8", errors="ignore", newline="") as fh:
                for row in csv.reader(fh):
                    lines.append("\t".join(row))
            return "\n".join(lines)

        except Exception as exc:
            logger.error("CSVExtractor failed for '%s': %s", path.name, exc)
            return ""


class PPTXExtractor(BaseExtractor):
    """
    Extracts text from PowerPoint presentations.

    Collects: slide titles, text boxes, table cells, and speaker notes.
    Each slide is prefixed with a slide-number header for readability.
    """

    _EXTENSIONS = frozenset({".pptx"})

    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def extract(self, path: Path) -> str:
        try:
            from pptx import Presentation  # noqa: PLC0415

            prs = Presentation(str(path))
            lines: list[str] = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                lines.append(f"--- Slide {slide_num} ---")

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in para.runs).strip()
                            if text:
                                lines.append(text)

                    # Extract table cell text
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = "\t".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip():
                                lines.append(row_text)

                # Speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        lines.append(f"[Notes]: {notes}")

            return "\n".join(lines)

        except Exception as exc:
            logger.error("PPTXExtractor failed for '%s': %s", path.name, exc)
            return ""


class XLSXExtractor(BaseExtractor):
    """
    Extracts cell content from Excel workbooks (.xlsx).

    Each sheet is prefixed with its name. Cells are tab-separated per row.
    Note: .xls (legacy format) requires the 'xlrd' library and is not
    handled here — use .xlsx where possible.
    """

    _EXTENSIONS = frozenset({".xlsx"})

    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def extract(self, path: Path) -> str:
        try:
            import openpyxl  # noqa: PLC0415

            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            lines: list[str] = []

            for sheet in wb.worksheets:
                lines.append(f"--- Sheet: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(
                        str(cell) for cell in row if cell is not None
                    )
                    if row_text.strip():
                        lines.append(row_text)

            wb.close()
            return "\n".join(lines)

        except Exception as exc:
            logger.error("XLSXExtractor failed for '%s': %s", path.name, exc)
            return ""


class ImageExtractor(BaseExtractor):
    """
    Extracts text from images using RapidOCR (OCR only — no captioning).

    Falls back to a short metadata description if RapidOCR is not
    installed or finds no text, so files remain indexable by name.
    """

    _EXTENSIONS = frozenset({".jpg", ".jpeg", ".png", ".bmp", ".gif", ".webp"})

    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def extract(self, path: Path) -> str:
        try:
            ocr_text = _run_ocr_on_bytes(path.read_bytes())
        except Exception as exc:
            logger.error("ImageExtractor failed for '%s': %s", path.name, exc)
            ocr_text = ""

        if ocr_text:
            return ocr_text

        # Fallback: metadata-only, so the file is still findable by name
        return f"Image file named {path.name} with extension {path.suffix.lower()}"


# ── Registry ──────────────────────────────────────────────────────────────────

def _build_registry() -> dict[str, BaseExtractor]:
    """Build a mapping of file extension → extractor instance."""
    registry: dict[str, BaseExtractor] = {}
    for extractor in [
        TxtExtractor(),
        HTMLExtractor(),
        PDFExtractor(),
        DOCXExtractor(),
        CSVExtractor(),
        PPTXExtractor(),
        XLSXExtractor(),
        ImageExtractor(),
    ]:
        for ext in extractor.supported_extensions():
            registry[ext] = extractor
    return registry


_REGISTRY: dict[str, BaseExtractor] = _build_registry()


# ── Metadata extraction ───────────────────────────────────────────────────────

def extract_metadata(path: Path) -> FileMetadata:
    """Read filesystem attributes and return a FileMetadata object."""
    stat = path.stat()
    return FileMetadata(
        name=path.name,
        path=str(path.resolve()),
        extension=path.suffix.lower(),
        size=stat.st_size,
        created_date=datetime.fromtimestamp(stat.st_ctime).isoformat(),
        modified_date=datetime.fromtimestamp(stat.st_mtime).isoformat(),
    )


# ── Text utilities ────────────────────────────────────────────────────────────

def normalize_text(text: str) -> str:
    """Collapse excess horizontal whitespace and blank lines."""
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def truncate_text(text: str, max_length: int = MAX_TEXT_LENGTH) -> str:
    """Truncate text to max_length at a word boundary, appending a notice."""
    if len(text) <= max_length:
        return text
    truncated = text[:max_length].rsplit(" ", 1)[0]
    return truncated + "\n\n[... text truncated for indexing ...]"


def chunk_text(
    text: str,
    size: int = DEFAULT_CHUNK_SIZE,
    overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[str]:
    """
    Split text into overlapping chunks for better retrieval quality.

    Chunk boundaries prefer sentence endings ('. ') then word boundaries.

    Args:
        text:    Text to split.
        size:    Target chunk size in characters (default 800).
        overlap: Characters from the previous chunk to repeat (default 150).

    Returns:
        List of non-empty text chunks.
    """
    if not text:
        return []

    chunks: list[str] = []
    start = 0

    while start < len(text):
        end = start + size
        chunk = text[start:end]

        # Break at sentence boundary if possible, else at word boundary
        if end < len(text):
            boundary = chunk.rfind(". ")
            if boundary == -1 or boundary < size // 2:
                boundary = chunk.rfind(" ")
            if boundary != -1:
                chunk = chunk[: boundary + 1]

        chunks.append(chunk.strip())
        
        if start + len(chunk) >= len(text):
            break
            
        start += max(1, len(chunk) - overlap)

    return [c for c in chunks if c]


# ── Public API ────────────────────────────────────────────────────────────────

def get_supported_extensions() -> set[str]:
    """Return all file extensions this module can process."""
    return set(_REGISTRY.keys())


def extract_text(file_path: str) -> str:
    """
    Extract plain text from a file (V1-compatible entry point).

    Routes to the appropriate extractor, normalises whitespace, and
    truncates to MAX_TEXT_LENGTH.

    Args:
        file_path: Absolute or relative path to the target file.

    Returns:
        Cleaned text string, or empty string on failure / unsupported type.
    """
    path = Path(file_path)

    if not path.exists():
        logger.warning("File not found: %s", path)
        return ""

    extractor = _REGISTRY.get(path.suffix.lower())
    if extractor is None:
        logger.warning("Unsupported file type '%s': %s", path.suffix, path.name)
        return ""

    raw = extractor.extract(path)
    return truncate_text(normalize_text(raw))


def extract_with_metadata(
    file_path: str,
    chunk: bool = True,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> Optional[ExtractionResult]:
    """
    Extract text, file metadata, and optional chunks from a file (V2 API).

    Args:
        file_path:     Path to the file.
        chunk:         Produce overlapping text chunks (default True).
        chunk_size:    Target chunk size in characters.
        chunk_overlap: Overlap between consecutive chunks.

    Returns:
        ExtractionResult on success, None if the file is missing or unsupported.
    """
    path = Path(file_path)

    if not path.exists():
        logger.warning("File not found: %s", path)
        return None

    extractor = _REGISTRY.get(path.suffix.lower())
    if extractor is None:
        logger.warning("Unsupported file type '%s': %s", path.suffix, path.name)
        return None

    raw = extractor.extract(path)
    content = truncate_text(normalize_text(raw))
    metadata = extract_metadata(path)
    chunks = chunk_text(content, chunk_size, chunk_overlap) if chunk else []

    return ExtractionResult(content=content, metadata=metadata, chunks=chunks)


# ── Dev entry point ───────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys

    target = sys.argv[1] if len(sys.argv) > 1 else "Sample_files/rrrrr.png"
    result = extract_with_metadata(target)

    if result:
        print(f"File      : {result.metadata.name}")
        print(f"Size      : {result.metadata.size} bytes")
        print(f"Modified  : {result.metadata.modified_date}")
        print(f"Chunks    : {len(result.chunks)}")
        print(f"\nContent preview:\n{result.content[:500]}")
    else:
        print("Extraction failed or unsupported file type.")